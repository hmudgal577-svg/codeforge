from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
PRIMARY = RGBColor(79, 70, 229)    # Indigo
SECONDARY = RGBColor(139, 92, 246) # Violet
DARK = RGBColor(30, 30, 50)
GRAY = RGBColor(100, 100, 120)
WHITE = RGBColor(255, 255, 255)
LIGHT_BG = RGBColor(248, 248, 255)
ACCENT = RGBColor(16, 185, 129)    # Green

IMG_DIR = os.path.join(os.path.dirname(__file__), 'img')

def add_bg(slide, color=LIGHT_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text(slide, left, top, w, h, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tf

def add_para(tf, text, size=16, bold=False, color=DARK, spacing=Pt(6)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.space_before = spacing
    return p

def add_image_safe(slide, name, left, top, w, h):
    path = os.path.join(IMG_DIR, name)
    if os.path.exists(path):
        slide.shapes.add_picture(path, left, top, w, h)

def add_bullet_slide(title, bullets, img_name=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    # Top bar
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PRIMARY)
    # Title
    add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7), title, size=32, bold=True, color=PRIMARY)
    # Divider
    add_shape(slide, Inches(0.8), Inches(1.1), Inches(2), Inches(0.05), SECONDARY)
    
    content_w = Inches(6.5) if img_name else Inches(11)
    tf = add_text(slide, Inches(0.8), Inches(1.5), content_w, Inches(5), "", size=16, color=DARK)
    tf.paragraphs[0].text = ""
    for b in bullets:
        if b.startswith("##"):
            add_para(tf, b[2:].strip(), size=18, bold=True, color=PRIMARY, spacing=Pt(14))
        else:
            add_para(tf, f"▸  {b}", size=15, color=GRAY, spacing=Pt(8))
    
    if img_name:
        add_image_safe(slide, img_name, Inches(7.8), Inches(1.5), Inches(4.8), Inches(4.8))
    
    # Slide number
    add_text(slide, Inches(12.3), Inches(7), Inches(0.8), Inches(0.4), str(len(prs.slides)), size=10, color=GRAY, align=PP_ALIGN.RIGHT)
    return slide

# ============ SLIDE 1: TITLE ============
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
add_shape(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), RGBColor(238, 238, 255))
add_shape(s, Inches(0), Inches(0), Inches(0.12), Inches(7.5), PRIMARY)
add_text(s, Inches(1.5), Inches(1.5), Inches(10), Inches(1), "CODEFORGE AI", size=52, bold=True, color=PRIMARY, align=PP_ALIGN.LEFT)
add_text(s, Inches(1.5), Inches(2.5), Inches(10), Inches(0.8), "AI-Powered Collaborative Cloud IDE", size=28, color=SECONDARY, align=PP_ALIGN.LEFT)
add_shape(s, Inches(1.5), Inches(3.5), Inches(3), Inches(0.05), ACCENT)
add_text(s, Inches(1.5), Inches(4), Inches(10), Inches(0.5), "A Full-Stack Development Environment with Real-Time Collaboration,", size=16, color=GRAY)
add_text(s, Inches(1.5), Inches(4.4), Inches(10), Inches(0.5), "Multi-Language Execution Engine & AI-Powered Code Intelligence", size=16, color=GRAY)
add_text(s, Inches(1.5), Inches(5.5), Inches(10), Inches(0.4), "Major Project Presentation  |  2025-26", size=14, color=GRAY)
add_text(s, Inches(1.5), Inches(6), Inches(10), Inches(0.4), "Department of Computer Science & Engineering", size=13, color=GRAY)
add_image_safe(s, 'hero.png', Inches(7.5), Inches(1.5), Inches(5), Inches(5))

# ============ SLIDE 2: INTRODUCTION ============
add_bullet_slide("Introduction", [
    "##What is CodeForge AI?",
    "A browser-based Cloud IDE built for modern developers",
    "Write, run, debug & collaborate on code — all in one place",
    "Supports 20+ programming languages with instant execution",
    "AI assistant built-in for code analysis & debugging",
    "##Why CodeForge?",
    "No setup needed — works entirely in the browser",
    "Real-time collaboration like Google Docs for code",
    "Intelligent AI assistant powered by multiple providers",
    "Secure sandboxed execution using Docker containers",
], 'hero.png')

# ============ SLIDE 3: PROBLEM STATEMENT ============
add_bullet_slide("Problem Statement", [
    "##Current Challenges in Development",
    "Setting up development environments is time-consuming",
    "Collaborating on code requires external tools (Git, Slack, etc.)",
    "Running multiple languages needs installing separate runtimes",
    "No integrated AI assistance in traditional IDEs",
    "Security risks when executing untrusted code",
    "##Our Solution",
    "Single platform for coding, execution, collaboration & AI",
    "Zero-setup browser-based IDE with Docker sandboxing",
    "Multi-provider AI failover ensures 99.9% availability",
])

# ============ SLIDE 4: OBJECTIVES ============
add_bullet_slide("Project Objectives", [
    "##Primary Objectives",
    "Build a cloud-based IDE accessible from any browser",
    "Support 20+ programming languages with sandboxed execution",
    "Integrate AI-powered code analysis with multi-provider failover",
    "Enable real-time multi-user collaboration with WebSocket",
    "##Secondary Objectives",
    "Implement security scanning before code execution",
    "Create Docker-based isolated runtime environments",
    "Build admin dashboard with analytics & monitoring",
    "Ensure responsive UI with modern design patterns",
])

# ============ SLIDE 5: TECHNOLOGY STACK ============
add_bullet_slide("Technology Stack", [
    "##Frontend",
    "Next.js 14 (React) — Server-side rendering & routing",
    "Monaco Editor — VS Code's editor engine",
    "TailwindCSS — Utility-first styling",
    "Socket.IO Client — Real-time communication",
    "##Backend",
    "NestJS (Node.js) — Enterprise-grade API framework",
    "Prisma ORM — Type-safe database access",
    "PostgreSQL — Primary relational database",
    "Redis — Caching & session management",
    "RabbitMQ — Message queue for async tasks",
], 'techstack.png')

# ============ SLIDE 6: SYSTEM ARCHITECTURE ============
add_bullet_slide("System Architecture", [
    "##Microservices Architecture",
    "API Gateway — Central routing & authentication",
    "Execution Engine — Sandboxed code runners (Docker)",
    "AI Service — Multi-provider failover (Gemini, Groq)",
    "Collaboration Gateway — WebSocket real-time sync",
    "##Infrastructure",
    "Docker containers for 20 language runtimes",
    "PostgreSQL for persistent data storage",
    "Redis for caching & real-time presence",
    "RabbitMQ for event-driven microservice communication",
], 'architecture.png')

# ============ SLIDE 7: CODE EDITOR ============
add_bullet_slide("Monaco Code Editor", [
    "##Features",
    "Syntax highlighting for 20+ languages",
    "IntelliSense auto-completion",
    "Multi-cursor editing support",
    "Theme customization (Dark/Light)",
    "File explorer with tree view",
    "##Integration",
    "Real-time content sync via WebSocket (CRDT)",
    "Auto-detect language from file extension",
    "Stdin input support for interactive programs",
    "One-click Run button with console output",
], 'hero.png')

# ============ SLIDE 8: MULTI-LANGUAGE EXECUTION ============
add_bullet_slide("Multi-Language Execution Engine", [
    "##20 Supported Languages",
    "JavaScript, TypeScript, Python, Java, C++, C",
    "Go, Rust, Ruby, PHP, Dart, Kotlin",
    "Scala, Swift, C#, Perl, R, Lua",
    "PowerShell, Bash",
    "##Execution Strategy",
    "Smart routing: Local runtime → Docker fallback",
    "Each language runs in isolated Docker container",
    "Memory limit: 128MB | CPU limit: 0.5 cores",
    "Automatic timeout: 30 seconds max",
    "Stdin/stdout piping for interactive programs",
])

# ============ SLIDE 9: DOCKER SANDBOXING ============
add_bullet_slide("Docker Sandboxed Execution", [
    "##Security Constraints",
    "Read-only filesystem (--read-only)",
    "No network access (--network=none)",
    "Non-root user execution (uid 1000)",
    "Process limit: 50 PIDs max (fork bomb prevention)",
    "No privilege escalation (no-new-privileges)",
    "##Container Architecture",
    "Alpine-based lightweight images (~50MB each)",
    "Separate /build directory for compiled languages",
    "Auto-cleanup after execution completes",
    "Parallel execution support for multiple users",
])

# ============ SLIDE 10: AI ASSISTANT ============
add_bullet_slide("AI-Powered Code Intelligence", [
    "##10 Analysis Types",
    "Code Explanation — Understand any code instantly",
    "Bug Detection — Find logical errors & edge cases",
    "Performance Optimization — Big-O & efficiency",
    "Security Vulnerability — OWASP Top 10 scanning",
    "Code Refactoring — SOLID principles & clean code",
    "##AI Chat Assistant",
    "Real-time chat within the IDE workspace",
    "Context-aware responses with conversation history",
    "Markdown rendering for formatted code output",
    "Multi-language support (Hindi & English)",
])

# ============ SLIDE 11: AI FAILOVER ============
add_bullet_slide("Multi-Provider AI Failover", [
    "##Provider Strategy",
    "Primary: Google Gemini 2.0 Flash (1500 req/day)",
    "Backup: Groq Llama 3.1 (14400 req/day)",
    "Tertiary: OpenRouter (200 req/day)",
    "##How Failover Works",
    "Automatic rotation when provider hits rate limit",
    "429 errors trigger instant switch to next provider",
    "Daily quota tracking per provider",
    "Smart fallback with mock responses if all fail",
    "99.9% AI availability guaranteed",
])

# ============ SLIDE 12: REAL-TIME COLLABORATION ============
add_bullet_slide("Real-Time Collaboration", [
    "##Features",
    "Multi-user cursor sync with color indicators",
    "Live presence bar showing online team members",
    "Real-time code editing (WebSocket + CRDT)",
    "Workspace chat with AI assistant integration",
    "##Technical Implementation",
    "Socket.IO for bi-directional real-time communication",
    "Room-based architecture for workspace isolation",
    "Conflict-free Replicated Data Types (CRDT) for sync",
    "30-second heartbeat for connection health",
], 'collab.png')

# ============ SLIDE 13: SECURITY ============
add_bullet_slide("Security Architecture", [
    "##Authentication & Authorization",
    "JWT + Refresh token based auth",
    "Bcrypt password hashing (12 salt rounds)",
    "Role-based access control (RBAC)",
    "##Code Execution Security",
    "Pre-execution threat scanning engine",
    "Pattern matching for dangerous code (fork bombs, etc.)",
    "Risk scoring system (0-100 scale)",
    "Auto-block for critical threats, warn for medium",
    "Docker container isolation with no network access",
], 'security.png')

# ============ SLIDE 14: DATABASE DESIGN ============
add_bullet_slide("Database Design", [
    "##PostgreSQL + Prisma ORM",
    "User — Authentication & profile data",
    "Workspace — Project containers with settings",
    "File — Source code files with versioning",
    "ExecutionJob — Code execution records & results",
    "AiRequest — AI analysis history & caching",
    "ChatMessage — Real-time collaboration messages",
    "Notification — User notification system",
    "##Key Design Decisions",
    "Prisma for type-safe queries & auto-migrations",
    "Connection pooling (17 connections) for performance",
    "Indexed queries for fast workspace lookups",
])

# ============ SLIDE 15: RESULTS & DEMO ============
add_bullet_slide("Results & Demo", [
    "##Key Metrics",
    "20 languages supported & tested successfully",
    "AI response time: < 3 seconds average",
    "Code execution: < 5 seconds for most languages",
    "Real-time sync latency: < 100ms",
    "##Live Demo Highlights",
    "C++ diamond pattern with stdin input (cin >> n)",
    "AI chat responding in Hindi & English",
    "Multi-user collaboration with live cursors",
    "Security engine blocking malicious code",
    "Docker fallback for uninstalled languages",
])

# ============ SLIDE 16: FUTURE SCOPE & THANK YOU ============
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
add_shape(s, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PRIMARY)
add_text(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7), "Future Scope & Conclusion", size=32, bold=True, color=PRIMARY)
add_shape(s, Inches(0.8), Inches(1.1), Inches(2), Inches(0.05), SECONDARY)

tf = add_text(s, Inches(0.8), Inches(1.5), Inches(5.5), Inches(4), "", size=16)
tf.paragraphs[0].text = ""
for item in [
    "##Future Enhancements",
    "OAuth login (Google, GitHub)",
    "Cloud deployment (Vercel + Railway)",
    "Git integration for version control",
    "Terminal emulator (xterm.js)",
    "Plugin marketplace for extensions",
    "Mobile responsive design",
    "AI auto-complete while typing",
]:
    if item.startswith("##"):
        add_para(tf, item[2:], size=18, bold=True, color=PRIMARY, spacing=Pt(14))
    else:
        add_para(tf, f"▸  {item}", size=14, color=GRAY, spacing=Pt(6))

# Thank you box
box = add_shape(s, Inches(7), Inches(1.5), Inches(5.5), Inches(4.5), RGBColor(238, 238, 255))
add_text(s, Inches(7.5), Inches(2), Inches(4.5), Inches(1), "Thank You!", size=42, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
add_text(s, Inches(7.5), Inches(3.2), Inches(4.5), Inches(0.5), "Questions & Discussion", size=20, color=SECONDARY, align=PP_ALIGN.CENTER)
add_shape(s, Inches(8.5), Inches(3.9), Inches(2.5), Inches(0.04), ACCENT)
add_text(s, Inches(7.5), Inches(4.3), Inches(4.5), Inches(0.4), "CodeForge AI — Major Project 2025-26", size=13, color=GRAY, align=PP_ALIGN.CENTER)
add_text(s, Inches(7.5), Inches(4.7), Inches(4.5), Inches(0.4), "Department of Computer Science & Engineering", size=12, color=GRAY, align=PP_ALIGN.CENTER)

# Save
output = os.path.join(os.path.dirname(__file__), 'CodeForge_AI_Presentation.pptx')
prs.save(output)
print(f"DONE! Saved: {output}")
